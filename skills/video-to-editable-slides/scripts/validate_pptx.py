#!/usr/bin/env python3
"""Validate basic PPTX package integrity and report slide statistics."""
import argparse, json, zipfile
from pathlib import Path
from pptx import Presentation

def main():
    ap = argparse.ArgumentParser(); ap.add_argument("pptx"); args = ap.parse_args()
    path = Path(args.pptx)
    if not path.is_file(): raise SystemExit(f"Missing file: {path}")
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip(); entries = set(zf.namelist())
        required = {"[Content_Types].xml", "ppt/presentation.xml"}
        missing = sorted(required - entries)
    prs = Presentation(path)
    slide_stats = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text.strip() for sh in slide.shapes if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        slide_stats.append({"slide": i, "shapes": len(slide.shapes), "text_chars": sum(map(len, texts))})
    report = {
        "file": str(path.resolve()), "bytes": path.stat().st_size,
        "zip_integrity": "ok" if bad is None else bad, "missing_required_entries": missing,
        "slides": len(prs.slides), "width_inches": prs.slide_width / 914400,
        "height_inches": prs.slide_height / 914400, "slide_stats": slide_stats
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    if bad or missing or not prs.slides: raise SystemExit(1)

if __name__ == "__main__": main()
