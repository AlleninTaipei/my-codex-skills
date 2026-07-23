#!/usr/bin/env python3
"""Build a clean editable 16:9 deck from the skill's slide-spec JSON."""
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def rgb(value):
    value = value.lstrip("#"); return RGBColor.from_string(value.upper())

def main():
    ap = argparse.ArgumentParser(); ap.add_argument("spec"); ap.add_argument("output"); args = ap.parse_args()
    data = json.loads(Path(args.spec).read_text(encoding="utf-8")); meta = data.get("metadata", {})
    theme = {"background":"#18181A","foreground":"#F2F0EB","accent":"#F49A19","muted":"#A5A5AA","font":"Microsoft JhengHei"}
    theme.update(meta.get("theme", {})); bg, fg, accent, muted = map(rgb, (theme["background"],theme["foreground"],theme["accent"],theme["muted"]))
    prs = Presentation(); prs.slide_width = Inches(13.333333); prs.slide_height = Inches(7.5)
    def text(slide,x,y,w,h,value,size=20,color=fg,bold=False,align=PP_ALIGN.LEFT):
        sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=str(value); r.font.name=theme["font"]; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; return sh
    def card(slide,x,y,w,h,heading,body,color):
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=rgb("#232327"); sh.line.color.rgb=color
        text(slide,x+.17,y+.12,w-.34,.4,heading,15,color,True); text(slide,x+.17,y+.56,w-.34,h-.68,body,12,fg)
    for idx, spec in enumerate(data.get("slides", []), 1):
        slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=bg
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.05)); bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
        layout=spec.get("layout","bullets"); title=spec.get("title",""); subtitle=spec.get("subtitle",""); section=spec.get("section","")
        if layout == "cover":
            text(slide,.8,1.25,11.7,1.2,title,42,fg,True); text(slide,.82,2.55,10.8,.65,subtitle,20,muted); text(slide,.82,6.7,11.5,.3,meta.get("source",""),9,muted)
            continue
        text(slide,.45,.18,5,.28,section,9,accent,True); text(slide,11.9,.18,.9,.28,f"{idx:02d}",9,muted,False,PP_ALIGN.RIGHT)
        text(slide,.75,.62,11.8,.72,title,28,fg,True)
        if subtitle: text(slide,.77,1.3,11.5,.4,subtitle,13,muted)
        elems=spec.get("elements",[])
        if layout in ("cards","comparison","process"):
            n=max(1,len(elems)); gap=.3; left=.75; total=11.83; w=(total-gap*(n-1))/n
            for i,e in enumerate(elems): card(slide,left+i*(w+gap),2.0,w,3.25,e.get("heading",e.get("text",f"{i+1}")),e.get("body",e.get("caption","")),rgb(e.get("accent",theme["accent"])))
        elif layout == "quote":
            q=elems[0].get("text","") if elems else ""; text(slide,1.1,2.05,11.1,2.4,q,30,accent,True,PP_ALIGN.CENTER)
        elif layout == "image":
            e=elems[0] if elems else {}; p=Path(e.get("path",""));
            if p.is_file(): slide.shapes.add_picture(str(p),Inches(1.0),Inches(1.85),width=Inches(11.3),height=Inches(4.8))
            if e.get("caption"): text(slide,1.0,6.72,11.3,.3,e["caption"],9,muted,False,PP_ALIGN.CENTER)
        else:
            lines=[e.get("text",e.get("body","")) for e in elems if e.get("text") or e.get("body")]
            sh=slide.shapes.add_textbox(Inches(.95),Inches(1.85),Inches(11.3),Inches(4.9)); tf=sh.text_frame; tf.clear(); tf.word_wrap=True
            for i,line in enumerate(lines):
                p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=line; p.font.name=theme["font"]; p.font.size=Pt(21); p.font.color.rgb=fg; p.space_after=Pt(12)
    prs.core_properties.title=meta.get("title",""); prs.core_properties.comments=f"Source: {meta.get('source','')}"; prs.save(args.output)
    print(f"saved={Path(args.output).resolve()} slides={len(prs.slides)}")

if __name__ == "__main__": main()
